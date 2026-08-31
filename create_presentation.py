"""Generate PowerPoint presentation for Нічний Дозор (current 4-script version)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(ROOT, "PresentationAssets", "Presentation_Nichniy_Dozor.pptx")
ASSETS = os.path.join(ROOT, "PresentationAssets")
SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(0x08, 0x0E, 0x18)
CARD = RGBColor(0x10, 0x18, 0x28)
ACCENT = RGBColor(0x4A, 0x9E, 0xFF)
GOLD = RGBColor(0xFF, 0xD7, 0x40)
GREEN = RGBColor(0x45, 0xC4, 0x70)
TEXT = RGBColor(0xEA, 0xF0, 0xFA)
DIM = RGBColor(0x8A, 0x98, 0xB0)

CHOSEN_MECHANICS = [
    ("21", "Стартовий вибір", "2 раси: Ельфи (+15% швидкість) / Гноми (+20% урон)"),
    ("27", "Часові обмеження", "Таймер на кожну хвилю — час вийшов = поразка"),
    ("11", "Захист цілі", "Кристал 100 HP — вороги йдуть по шляху до нього"),
]


def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = BG


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    return s


def center_card(slide, w, h, y=None):
    if y is None:
        y = (SLIDE_H - h) / 2
    x = (SLIDE_W - w) / 2
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = CARD
    sh.line.color.rgb = ACCENT
    sh.line.width = Pt(1.5)


def title(slide, text, subtitle=None, y=0.42):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(SLIDE_W - 1.4), Inches(0.85))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(1.0), Inches(y + 0.62), Inches(SLIDE_W - 2.0), Inches(0.45))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.color.rgb = DIM
        sp.alignment = PP_ALIGN.CENTER


def bullets(slide, items, top=1.45, width=10.5, size=19, center=True):
    x = (SLIDE_W - width) / 2
    tb = slide.shapes.add_textbox(Inches(x), Inches(top), Inches(width), Inches(SLIDE_H - top - 0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT if line else DIM
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        p.space_after = Pt(8)


def slide_title(prs):
    s = blank(prs)
    center_card(s, 9.8, 4.0, 1.75)
    tb = s.shapes.add_textbox(Inches(2.0), Inches(2.55), Inches(9.3), Inches(2.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "НІЧНИЙ ДОЗОР"
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Tower Defense · Unity · URP"
    p2.font.size = Pt(22)
    p2.font.color.rgb = GOLD
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    p3 = tf.add_paragraph()
    p3.text = "Захист кристала вночі — 5 хвиль, 2 башні, 2 раси"
    p3.font.size = Pt(17)
    p3.font.color.rgb = DIM
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(16)


def slide_random(prs):
    s = blank(prs)
    title(s, "1. Механіки з RANDOM.ORG", "Обрано по 1 механіці з кожного ряду")
    center_card(s, 11.0, 5.3, 1.05)
    y = 1.55
    for num, name, desc in CHOSEN_MECHANICS:
        tb = s.shapes.add_textbox(Inches(1.8), Inches(y), Inches(9.8), Inches(1.1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = f"№{num} · {name}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        y += 1.55


def slide_game_overview(prs):
    s = blank(prs)
    title(s, "2. Що в грі зараз")
    center_card(s, 10.8, 5.2, 1.15)
    bullets(s, [
        "2 раси · 2 типи башен · 3 типи ворогів",
        "5 хвиль з таймером на кожну",
        "Стартове золото: 120 · Кристал: 100 HP",
        "Лучник (50) — ближній, швидкий  |  Гармата (90) — дальня",
        "Карта збережена в сцені SampleScene",
        "Код спрощено до 4 скриптів",
    ], top=1.55, size=18)


def slide_architecture(prs):
    s = blank(prs)
    title(s, "3. Архітектура коду", "Assets/Scripts/")
    center_card(s, 11.2, 5.4, 1.0)

    rows = [
        ("GameConfig.cs", "Баланс, хвилі, шлях ворогів"),
        ("BuildZone.cs", "Зелена клітинка на сцені"),
        ("Game.cs", "Логіка + Tower + Enemy"),
        ("UIManager.cs", "Меню, HUD, кнопки"),
    ]
    y = 1.45
    for file, role in rows:
        tb = s.shapes.add_textbox(Inches(1.6), Inches(y), Inches(10.2), Inches(0.75))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = file
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = role
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT
        p2.alignment = PP_ALIGN.CENTER
        y += 1.15

    tb = s.shapes.add_textbox(Inches(1.5), Inches(6.35), Inches(10.3), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Сцена: Level (карта) + NightWatch (Game, UIManager) + Main Camera"
    p.font.size = Pt(14)
    p.font.color.rgb = DIM
    p.alignment = PP_ALIGN.CENTER


def slide_flow(prs):
    s = blank(prs)
    title(s, "4. Потік гри")
    center_card(s, 10.5, 5.3, 1.05)
    bullets(s, [
        "Play → меню: обрати расу",
        "HUD → обрати башню (Лучник / Гармата)",
        "Клік по зеленій клітинці → побудова",
        "Кнопка «Хвиля» → спавн ворогів + таймер",
        "Башні б'ють найближчого ворога (без снарядів)",
        "Усі вороги мертві до кінця часу → наступна хвиля",
        "5 хвиль → перемога  |  час або кристал → поразка",
    ], top=1.45, size=17)


def slide_game_class(prs):
    s = blank(prs)
    title(s, "5. Клас Game — «мозок» гри")
    center_card(s, 11.0, 5.3, 1.05)
    bullets(s, [
        "Game.I — singleton (один менеджер на гру)",
        "Start: знайти Level, зібрати BuildZone, показати меню",
        "Update: таймер хвилі, HUD, кліки по карті",
        "StartWithRace / SelectTower / StartNextWave",
        "SpawnWave — корутина з паузою між ворогами",
        "Tower.Create — миттєвий урон  |  Enemy — рух по EnemyPath",
    ], top=1.5, size=17)


def slide_config(prs):
    s = blank(prs)
    title(s, "6. GameConfig.cs")
    center_card(s, 10.8, 5.2, 1.15)
    bullets(s, [
        "static class — усі числа в одному місці",
        "Масив Waves: скільки ворогів на кожну хвилю",
        "Масив EnemyPath: точки маршруту",
        "TowerCost, TowerRange, TowerDamage, TowerFireRate",
        "ApplyRace: ельфи швидше, гноми сильніше",
        "WaveSeconds: 45 + wave×5 секунд",
    ], top=1.55, size=17)


def slide_ui(prs):
    s = blank(prs)
    title(s, "7. UIManager.cs")
    center_card(s, 10.5, 5.0, 1.2)
    bullets(s, [
        "UI створюється в коді (Canvas, кнопки, тексти)",
        "3 екрани: меню · HUD · кінець гри",
        "Ціна башен над кнопками (50 / 90 зол.)",
        "Підсвітка обраної башні зеленим",
        "RefreshHud — золото, HP, таймер, хвиля",
    ], top=1.6, size=18)


def slide_controls(prs):
    s = blank(prs)
    title(s, "8. Керування")
    center_card(s, 9.5, 4.5, 1.45)
    bullets(s, [
        "Ельфи / Гноми — вибір раси",
        "Лучник / Гармата — тип башні",
        "Клік по зеленій зоні — побудова",
        "Хвиля — старт наступної хвилі",
        "Знову — рестарт після перемоги/поразки",
    ], top=1.75, size=20)


def slide_exam(prs):
    s = blank(prs)
    title(s, "9. Питання на захисті", "Короткі відповіді")
    center_card(s, 11.2, 5.4, 1.0)
    bullets(s, [
        "Скільки скриптів? — 4 (GameConfig, BuildZone, Game, UIManager)",
        "Де карта? — об'єкт Level на сцені",
        "Як вороги йдуть? — по масиву EnemyPath",
        "Що таке корутина? — SpawnWave з yield return WaitForSeconds",
        "Коли поразка? — час вийшов або кристал знищено",
        "Документація: PresentationAssets/exam_prep_code_qa.pdf",
    ], top=1.4, size=16)


def slide_demo(prs):
    s = blank(prs)
    title(s, "10. Демонстрація")
    center_card(s, 10, 4.5, 1.5)
    tb = s.shapes.add_textbox(Inches(2.2), Inches(3.0), Inches(9), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "▶  LIVE DEMO / ВІДЕО"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Unity → SampleScene → Play"
    p2.font.size = Pt(18)
    p2.font.color.rgb = DIM
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)


def slide_thanks(prs):
    s = blank(prs)
    center_card(s, 8, 2.5, 2.5)
    tb = s.shapes.add_textbox(Inches(2.8), Inches(2.9), Inches(7.8), Inches(1.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Дякуємо за увагу!"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Нічний Дозор · GitHub: Ekzamen_Bozhko"
    p2.font.size = Pt(18)
    p2.font.color.rgb = DIM
    p2.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide_title(prs)
    slide_random(prs)
    slide_game_overview(prs)
    slide_architecture(prs)
    slide_flow(prs)
    slide_game_class(prs)
    slide_config(prs)
    slide_ui(prs)
    slide_controls(prs)
    slide_exam(prs)
    slide_demo(prs)
    slide_thanks(prs)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    try:
        prs.save(OUT)
        print(f"Saved: {OUT}")
    except PermissionError:
        alt = OUT.replace(".pptx", "_new.pptx")
        prs.save(alt)
        print(f"File locked — saved: {alt}")


if __name__ == "__main__":
    main()
